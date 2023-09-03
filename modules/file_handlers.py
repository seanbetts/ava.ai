# modules/file_handlers.py

import chainlit as cl

import io
import PyPDF2
from pptx import Presentation
from docx import Document
import pandas as pd
import json

from .utils import (extract_first_200_words, generate_actions, send_file_message, num_tokens_from_string, get_token_limit, is_over_token_limit, dataframe_to_json_metadata)

##--HANDLE TEXT--##
async def handle_text_file(uploaded_file, model):
    text = uploaded_file.content.decode("utf-8")
    extracted_text = extract_first_200_words(text)
    action_keys = ["question", "summarise", "bulletpoint_summary", "create_wordcloud", "get_themes", "get_quotes", "copy", "save_to_knowledgebase"]
    await send_file_message(uploaded_file.name, text, extracted_text, model, action_keys)

##--HANDLE DOCS--##
async def handle_doc_file(uploaded_file, model):
    doc_content = io.BytesIO(uploaded_file.content)
    doc = Document(doc_content)
    doc_text = "\n".join([para.text for para in doc.paragraphs])
    extracted_text = extract_first_200_words(doc_text)
    action_keys = ["question", "summarise", "bulletpoint_summary", "create_wordcloud", "get_themes", "get_quotes", "copy", "save_to_knowledgebase"]
    await send_file_message(uploaded_file.name, doc_text, extracted_text, model, action_keys)

##--HANDLE PDFS--##
async def handle_pdf_file(uploaded_file, model):
    pdf_content = io.BytesIO(uploaded_file.content)
    pdf_reader = PyPDF2.PdfReader(pdf_content)
    pdf_text = "\n".join([page.extract_text() for page in pdf_reader.pages])
    extracted_text = extract_first_200_words(pdf_text)
    action_keys = ["question", "summarise", "bulletpoint_summary", "create_wordcloud", "get_themes", "get_quotes", "copy", "save_to_knowledgebase", "upload_file"]
    await send_file_message(uploaded_file.name, pdf_text, extracted_text, model, action_keys)

##--HANDLE PPTS--##
async def handle_ppt_file(uploaded_file, model):
    ppt_content = io.BytesIO(uploaded_file.content)
    presentation = Presentation(ppt_content)
    ppt_text = "\n".join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])
    extracted_text = extract_first_200_words(ppt_text)
    action_keys = ["question", "summarise", "bulletpoint_summary", "create_wordcloud", "get_themes", "get_quotes", "copy", "save_to_knowledgebase", "upload_file"]
    await send_file_message(uploaded_file.name, ppt_text, extracted_text, model, action_keys)

###--HANDLE XLSXs--###
async def handle_xlsx_file(uploaded_file, model):
    data = io.BytesIO(uploaded_file.content)
    df = pd.read_excel(data)

    # Convert first 5 rows to markdown
    markdown_content = df.head().to_markdown()

    # Store the dataframe in the user session
    cl.user_session.set("df", df)

    text_file=dataframe_to_json_metadata(df)
    text_file_string = json.dumps(text_file)

    tokens = num_tokens_from_string(text_file_string, model)
    token_limit = get_token_limit(model)
    is_over = is_over_token_limit(tokens, token_limit)

    elements = [
        cl.Text(name="Here are the top 5 rows of data:", content=markdown_content, display="inline")
    ]

    action_keys = ["question", "get_insights", "upload_file"]
    actions = generate_actions("data", action_keys)

    await cl.Message(
        content=f"`{uploaded_file.name}` uploaded.\n It has {format(df.shape[0], ',')} rows and {format(df.shape[1], ',')} columns which is c.{format(tokens, ',')} tokens.\nYou're currently using the {model} model which has a token limit of {format(token_limit, ',')}.\n{is_over[1]}", elements=elements, actions=actions
    ).send()

###--HANDLE CSVs--###
async def handle_csv_file(uploaded_file, model):
    data = io.BytesIO(uploaded_file.content)
    df = pd.read_csv(data)

    # Convert first 5 rows to markdown
    markdown_content = df.head().to_markdown()

    # Store the dataframe in the user session
    cl.user_session.set("df", df)

    text_file=dataframe_to_json_metadata(df)
    text_file_string = json.dumps(text_file)

    tokens = num_tokens_from_string(text_file_string, model)
    token_limit = get_token_limit(model)
    is_over = is_over_token_limit(tokens, token_limit)

    elements = [
        cl.Text(name="Here are the top 5 rows of data:", content=markdown_content, display="inline")
    ]

    action_keys = ["question", "get_insights", "upload_file"]
    actions = generate_actions("data", action_keys)

    await cl.Message(
        content=f"`{uploaded_file.name}` uploaded.\n It has {format(df.shape[0], ',')} rows and {format(df.shape[1], ',')} columns which is c.{format(tokens, ',')} tokens.\nYou're currently using the {model} model which has a token limit of {format(token_limit, ',')}.\n{is_over[1]}", elements=elements, actions=actions
    ).send()

###--HANDLE IMAGEs--###
async def handle_image_file(uploaded_file):
    image_data = uploaded_file.content

    # Sending an image with the file name
    elements = [
    cl.Image(name=uploaded_file.name, display="inline", content=image_data)
    ]

    await cl.Message(content="Here's your image:", elements=elements).send()